import boto3
import json
import os
import pickle
import numpy as np
import faiss
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

# AWS Clients
textract = boto3.client(
    'textract',
    region_name='us-east-1',
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

bedrock = boto3.client(
    'bedrock-runtime',
    region_name='us-east-1',
    aws_access_key_id=os.getenv('BEDROCK_ACCESS_KEY'),
    aws_secret_access_key=os.getenv('BEDROCK_SECRET_KEY')
)

# FAISS index and metadata
EMBEDDING_DIM = 512
INDEX_PATH = 'doc_index.faiss'
META_PATH = 'doc_meta.pkl'

doc_index = faiss.IndexFlatIP(EMBEDDING_DIM)
doc_meta = []  # list of {path, filename, snippet}

def extract_text_local(file_path: str) -> str:
    """Extract text from PDF, DOCX, PPTX, XLSX, TXT and more."""
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        if ext == '.pdf':
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            text = " ".join(page.get_text() for page in doc)
            doc.close()
            return text

        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(file_path)
            return " ".join(para.text for para in doc.paragraphs)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                texts.append(cell.text)
            return " ".join(texts)

        elif ext in ('.xlsx', '.xls'):
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return " ".join(texts)

        elif ext in ('.txt', '.md', '.csv', '.log', '.json', '.xml', '.html'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        else:
            return ""
    except Exception as e:
        print(f"⚠️ Local extraction failed for {file_path}: {e}")
        return ""

def extract_text_textract(file_path: str) -> str:
    """Use AWS Textract for scanned PDFs and images"""
    try:
        with open(file_path, 'rb') as f:
            file_bytes = f.read()

        response = textract.detect_document_text(
            Document={'Bytes': file_bytes}
        )

        text = " ".join(
            block['Text']
            for block in response['Blocks']
            if block['BlockType'] == 'LINE'
        )
        return text
    except Exception as e:
        print(f"⚠️ Textract failed for {file_path}: {e}")
        return ""

def get_embedding(text: str) -> list:
    """Convert text to 512-dim vector using Titan Embeddings V2"""
    try:
        text = text[:8000]

        response = bedrock.invoke_model(
            modelId=os.getenv('TITAN_MODEL_ID', 'amazon.titan-embed-text-v2:0'),
            body=json.dumps({
                "inputText": text,
                "dimensions": 512,
                "normalize": True
            })
        )
        result = json.loads(response['body'].read())
        return result['embedding']
    except Exception as e:
        print(f"⚠️ Embedding failed: {e}")
        return None

def index_single_doc(file_path: str):
    """Index a single document"""
    global doc_index, doc_meta

    print(f"Indexing: {file_path}")

    # Step 1: Extract text
    text = extract_text_local(file_path)

    # Step 2: Fallback to Textract if local extraction empty
    if len(text.strip()) < 50:
        print(f"  → Using Textract for: {file_path}")
        text = extract_text_textract(file_path)

    if not text.strip():
        print(f"  ⚠️ No text found in: {file_path}")
        return False

    # Step 3: Get embedding
    embedding = get_embedding(text)
    if embedding is None:
        return False

    # Step 4: Add to FAISS index
    vec = np.array([embedding], dtype='float32')
    doc_index.add(vec)

    # Step 5: Store metadata
    snippet = text[:300].strip()
    doc_meta.append({
        "path": file_path,
        "filename": Path(file_path).name,
        "snippet": snippet,
        "full_text": text[:2000]
    })

    print(f"  ✅ Indexed: {Path(file_path).name}")
    return True

def index_docs_folder(folder_path: str):
    """Index all documents in a folder"""
    global doc_index, doc_meta

    folder = Path(folder_path)
    doc_extensions = {'.pdf', '.docx', '.doc', '.txt', '.md', '.csv', '.pptx', '.ppt', '.xlsx', '.xls'}

    total = 0
    success = 0

    for file_path in folder.rglob('*'):
        if file_path.suffix.lower() in doc_extensions:
            total += 1
            if index_single_doc(str(file_path)):
                success += 1

    save_index()

    print(f"\n✅ Indexing complete! {success}/{total} documents indexed")
    return {"total": total, "indexed": success}

def save_index():
    """Save FAISS index and metadata to disk"""
    faiss.write_index(doc_index, INDEX_PATH)
    with open(META_PATH, 'wb') as f:
        pickle.dump(doc_meta, f)
    print(f"💾 Index saved! ({len(doc_meta)} documents)")

def load_index():
    """Load FAISS index from disk"""
    global doc_index, doc_meta
    try:
        if os.path.exists(INDEX_PATH) and os.path.exists(META_PATH):
            doc_index = faiss.read_index(INDEX_PATH)
            with open(META_PATH, 'rb') as f:
                doc_meta = pickle.load(f)
            print(f"✅ Index loaded: {len(doc_meta)} documents")
            return True
    except Exception as e:
        print(f"⚠️ Could not load index: {e}")
    return False

# Test it
if __name__ == "__main__":
    print("Testing doc_indexer...")
    emb = get_embedding("This is a test document about project budgets")
    if emb:
        print(f"✅ Embedding works: {len(emb)}-dim vector")
    else:
        print("❌ Embedding failed")